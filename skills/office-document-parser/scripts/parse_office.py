import os
import sys
import subprocess
import argparse

def install_deps():
    deps = ['pandas', 'openpyxl', 'python-docx', 'python-pptx', 'tabulate', 'xlrd']
    missing = False
    try:
        import pandas
        import openpyxl
        import docx
        import pptx
        import tabulate
        import xlrd
    except ImportError:
        missing = True

    if missing:
        print("Missing dependencies, installing pandas, openpyxl, python-docx, python-pptx, tabulate...")
        subprocess.check_call([sys.executable, '-m', 'pip', 'install', *deps])
        print("Dependencies installed successfully.")

def parse_docx(file_path):
    import docx
    doc = docx.Document(file_path)
    md_lines = []
    md_lines.append(f"# Word Document: {os.path.basename(file_path)}\n")
    for para in doc.paragraphs:
        if para.text.strip():
            if para.style.name.startswith('Heading'):
                level_str = para.style.name.replace('Heading', '').strip()
                level = int(level_str) if level_str.isdigit() else 1
                level = min(level, 6)
                md_lines.append(f"{'#' * level} {para.text.strip()}")
            else:
                md_lines.append(para.text.strip())
            md_lines.append("")

    if doc.tables:
        md_lines.append("## Tables\n")
        for i, table in enumerate(doc.tables):
            md_lines.append(f"### Table {i+1}")
            for r_idx, row in enumerate(table.rows):
                row_data = [cell.text.replace("\n", " ").replace("\r", " ").strip() for cell in row.cells]
                md_lines.append("| " + " | ".join(row_data) + " |")
                if r_idx == 0:
                    md_lines.append("|" + "|".join(["---"] * len(row.cells)) + "|")
            md_lines.append("")
    return "\n".join(md_lines)

def parse_xlsx(file_path):
    import pandas as pd
    md_lines = []
    md_lines.append(f"# Excel Document: {os.path.basename(file_path)}\n")
    excel_data = pd.read_excel(file_path, sheet_name=None, dtype=str)
    for sheet_name, df in excel_data.items():
        md_lines.append(f"## Sheet: {sheet_name}\n")
        df = df.fillna("")
        if not df.empty:
            md_lines.append(df.to_markdown(index=False))
        else:
            md_lines.append("*Empty Sheet*")
        md_lines.append("\n")
    return "\n".join(md_lines)

def parse_pptx(file_path):
    import pptx
    prs = pptx.Presentation(file_path)
    md_lines = []
    md_lines.append(f"# PowerPoint Document: {os.path.basename(file_path)}\n")
    for i, slide in enumerate(prs.slides):
        md_lines.append(f"## Slide {i+1}")
        
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                md_lines.append(f"> **Notes**: {notes}\n")

        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                md_lines.append(shape.text.strip() + "\n")
            if shape.has_table:
                table = shape.table
                for r_idx, row in enumerate(table.rows):
                    row_data = [cell.text_frame.text.replace("\n", " ").replace("\r", " ").strip() for cell in row.cells]
                    md_lines.append("| " + " | ".join(row_data) + " |")
                    if r_idx == 0:
                        md_lines.append("|" + "|".join(["---"] * len(row.cells)) + "|")
                md_lines.append("\n")
    return "\n".join(md_lines)

def main():
    parser = argparse.ArgumentParser(description="Parse Office documents into Markdown.")
    parser.add_argument("file_path", help="Path to the Office file (.docx, .xlsx, .pptx)")
    args = parser.parse_args()

    install_deps()

    file_path = args.file_path
    if not os.path.exists(file_path):
        print(f"Error: File '{file_path}' does not exist.")
        sys.exit(1)

    ext = os.path.splitext(file_path)[1].lower()
    
    try:
        if ext == ".docx":
            output = parse_docx(file_path)
        elif ext in [".xlsx", ".xls", ".xlsm"]:
            output = parse_xlsx(file_path)
        elif ext == ".pptx":
            output = parse_pptx(file_path)
        else:
            print(f"Error: Unsupported file extension '{ext}'. Only .docx, .xlsx, .pptx are supported.")
            sys.exit(1)

        if hasattr(sys.stdout, 'reconfigure'):
            sys.stdout.reconfigure(encoding='utf-8')
        print(output)
    except Exception as e:
        print(f"Error parsing file: {e}")
        sys.exit(1)

if __name__ == "__main__":
    main()
